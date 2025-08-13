"""
PPT generation service using python-pptx

This service handles the generation of PowerPoint presentations from layout and content data.
The generated presentations follow a responsive design approach with proper positioning,
theming, and content formatting that can be customized through the frontend interface.
"""

import logging
from typing import Dict, List, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

class PPTService:
    """
    Service for generating PowerPoint presentations
    
    This class handles the core functionality for converting frontend layout
    configurations and AI-generated content into professional PowerPoint files.
    The service supports multiple themes and section types that can be controlled
    through the frontend interface.
    """
    
    def __init__(self):
        # Standard 16:9 aspect ratio dimensions for modern presentations
        self.slide_width = 13.33  # Inches - matches widescreen format
        self.slide_height = 7.5   # Inches - maintains 16:9 ratio
        
    async def generate_ppt_from_layout(
        self,
        layout: Dict[str, Any],
        content: Dict[str, Any],
        output_path: str,
        theme: str = "default"
    ) -> str:
        """
        Generate a PowerPoint file from layout and content
        
        This is the main entry point called by the frontend. It takes the layout
        configuration from the builder interface and the AI-generated content
        to create a complete PowerPoint presentation.
        
        Args:
            layout: Slide layout structure from frontend builder
                   Contains sections array with type, position, and styling info
            content: Generated content for each section from AI service
                    Keyed by section_0, section_1, etc. with content property
            output_path: Path to save the PPT file (handled by backend file service)
            theme: Slide theme selected in frontend theme picker
                  Options: "professional", "creative", "minimal", "colorful", "default"
            
        Returns:
            Path to the generated PPT file for download
        """
        try:
            # Create a new presentation with blank template
            prs = Presentation()
            
            # Set slide size to 16:9 aspect ratio (modern presentation standard)
            prs.slide_width = Inches(self.slide_width)
            prs.slide_height = Inches(self.slide_height)
            
            # Add a single slide with blank layout (index 6 = no predefined elements)
            slide_layout = prs.slide_layouts[6]  # Blank layout for full customization
            slide = prs.slides.add_slide(slide_layout)
            
            # Apply the selected theme styling (background, colors, gradients)
            self._apply_theme(slide, theme)
            
            # Add title if present in layout configuration
            if "title" in layout:
                self._add_title(slide, layout["title"], theme)
            
            # Process each section from the frontend builder
            sections = layout.get("sections", [])
            for i, section in enumerate(sections):
                # Get corresponding content for this section
                section_content = content.get(f"section_{i}", {})
                logger.info(f"Adding section {i}: {section.get('type', 'unknown')} with content: {section_content.get('content', 'No content')[:100]}...")
                # Add the section with positioning and styling
                self._add_section(slide, section, section_content, theme)
            
            # Save the presentation to the specified path
            prs.save(output_path)
            logger.info(f"PPT file generated: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating PPT file: {e}")
            # Fallback: create a minimal error slide for user feedback
            return self._create_error_ppt(output_path, str(e))
    
    def _apply_theme(self, slide, theme: str):
        """
        Apply theme styling to the slide background
        
        This method applies the visual theme selected in the frontend theme picker.
        Each theme has its own color palette and styling approach to match different
        presentation contexts (professional meetings, creative pitches, etc.).
        
        Available themes correspond to frontend theme options:
        - professional: Corporate blue gradient
        - creative: Modern purple gradient  
        - minimal: Clean white background
        - colorful: Vibrant red-teal gradient
        - default: Standard purple gradient
        """
        try:
            # Get background shape for styling
            background = slide.background
            fill = background.fill
            
            if theme == "professional":
                # Professional blue gradient - suitable for corporate presentations
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor(44, 62, 80)   # Dark blue-gray
                fill.gradient_stops[1].color.rgb = RGBColor(52, 73, 94)   # Medium blue-gray
            elif theme == "creative":
                # Creative purple gradient - modern and engaging for creative content
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor(102, 126, 234) # Light purple
                fill.gradient_stops[1].color.rgb = RGBColor(118, 75, 162)  # Deep purple
            elif theme == "minimal":
                # Minimal white background - clean and distraction-free
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255) # Pure white
            elif theme == "colorful":
                # Colorful gradient - vibrant and attention-grabbing
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor(255, 107, 107) # Coral red
                fill.gradient_stops[1].color.rgb = RGBColor(78, 205, 196)  # Teal
            else:
                # Default gradient - balanced purple theme
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor(102, 126, 234) # Light purple
                fill.gradient_stops[1].color.rgb = RGBColor(118, 75, 162)  # Deep purple
                
        except Exception as e:
            logger.warning(f"Error applying theme {theme}: {e}")
    
    def _add_title(self, slide, title: str, theme: str):
        """
        Add title to the slide with theme-appropriate styling
        
        Creates a prominent title at the top of the slide using the title text
        provided by the frontend. The title is styled to be highly visible and
        properly contrasted against the selected theme background.
        """
        try:
            # Position title at top of slide with generous margins
            left = Inches(0.5)      # Small left margin
            top = Inches(0.3)       # Small top margin  
            width = Inches(self.slide_width - 1)  # Full width minus margins
            height = Inches(1.2)    # Adequate height for large title text
            
            # Create text box for title
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title
            
            # Center-align the title for professional appearance
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Set large, bold font for maximum impact
            font = p.font
            font.size = Pt(32)  # Large title size
            font.bold = True    # Bold weight for emphasis
            
            # Choose text color based on theme background
            if theme == "minimal":
                font.color.rgb = RGBColor(51, 51, 51)    # Dark gray on white
            else:
                font.color.rgb = RGBColor(255, 255, 255) # White on colored backgrounds
            
            logger.info(f"Added title: {title}")
                
        except Exception as e:
            logger.error(f"Error adding title: {e}")
    
    def _add_section(self, slide, section: Dict[str, Any], content: Dict[str, Any], theme: str):
        """
        Add a section to the slide based on frontend layout configuration
        
        This method processes each section from the frontend builder, which includes:
        - Section type (determines content formatting)
        - Position data (x, y coordinates and dimensions as percentages)
        - Content data (text, bullet points, etc. from AI generation)
        
        Position data uses percentage-based coordinates that are converted to
        absolute positioning on the slide canvas.
        """
        try:
            # Get section type from frontend configuration
            section_type = section.get("type", "text")
            # Get position data (percentages from frontend drag-and-drop interface)
            position = section.get("position", {"x": 10, "y": 30, "width": 80, "height": 60})
            
            # Validate and clamp position values to prevent off-slide content
            x = max(0, min(95, position.get("x", 10)))      # X position: 0-95%
            y = max(0, min(95, position.get("y", 30)))      # Y position: 0-95%
            width = max(5, min(95, position.get("width", 80)))   # Width: 5-95%
            height = max(5, min(95, position.get("height", 60))) # Height: 5-95%
            
            # Ensure section doesn't overflow slide boundaries
            if x + width > 95:
                width = 95 - x
            if y + height > 95:
                height = 95 - y
            
            # Convert percentage coordinates to absolute inches
            left = Inches((x / 100) * self.slide_width)
            top = Inches((y / 100) * self.slide_height)
            width_inches = Inches((width / 100) * self.slide_width)
            height_inches = Inches((height / 100) * self.slide_height)
            
            logger.info(f"Adding section '{section_type}' at position: x={x}%, y={y}%, w={width}%, h={height}%")
            
            # Route to appropriate section handler based on type
            if section_type == "text":
                self._add_text_section(slide, content, left, top, width_inches, height_inches, theme)
            elif section_type == "bullet_list":
                self._add_bullet_list_section(slide, content, left, top, width_inches, height_inches, theme)
            elif section_type == "image":
                self._add_image_section(slide, content, left, top, width_inches, height_inches, theme)
            elif section_type == "highlight_box":
                self._add_highlight_box_section(slide, content, left, top, width_inches, height_inches, theme)
            elif section_type == "quote":
                self._add_quote_section(slide, content, left, top, width_inches, height_inches, theme)
            else:
                # Fallback to text for unknown section types
                self._add_text_section(slide, content, left, top, width_inches, height_inches, theme)
                
        except Exception as e:
            logger.error(f"Error adding section: {e}")
    
    def _add_text_section(self, slide, content: Dict[str, Any], left, top, width, height, theme: str):
        """
        Add a standard text section to the slide
        
        Creates a text box with AI-generated content formatted for readability.
        This is the most common section type used for paragraphs, descriptions,
        and general content blocks.
        """
        try:
            # Create text box at specified position and size
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Insert AI-generated content with fallback message
            section_content = content.get("content", "Content not available")
            tf.text = section_content
            
            # Apply readable text styling
            p = tf.paragraphs[0]
            font = p.font
            font.size = Pt(16)  # Medium size for body text readability
            
            # Set text color based on theme for proper contrast
            if theme == "minimal":
                font.color.rgb = RGBColor(51, 51, 51)    # Dark gray on white
            else:
                font.color.rgb = RGBColor(255, 255, 255) # White on colored backgrounds
            
            logger.info(f"Added text section with content: {section_content[:100]}...")
                
        except Exception as e:
            logger.error(f"Error adding text section: {e}")
    
    def _add_bullet_list_section(self, slide, content: Dict[str, Any], left, top, width, height, theme: str):
        """
        Add a bullet list section to the slide
        
        Formats AI-generated content as a bulleted list by splitting on newlines.
        Each line becomes a separate bullet point for better content organization
        and visual hierarchy.
        """
        try:
            # Create text box for bullet list
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Get content and split into individual bullet points
            section_content = content.get("content", "")
            lines = section_content.split('\n')  # Split on newlines for bullets
            
            # Add each line as a separate paragraph (bullet point)
            for i, line in enumerate(lines):
                if i == 0:
                    # First line uses existing paragraph
                    tf.text = line
                    p = tf.paragraphs[0]
                else:
                    # Additional lines create new paragraphs
                    p = tf.add_paragraph()
                    p.text = line
                
                # Style each bullet point for readability
                font = p.font
                font.size = Pt(14)  # Slightly smaller than body text
                
                # Apply theme-appropriate text color
                if theme == "minimal":
                    font.color.rgb = RGBColor(51, 51, 51)    # Dark gray
                else:
                    font.color.rgb = RGBColor(255, 255, 255) # White
            
            logger.info(f"Added bullet list section with {len(lines)} lines")
                    
        except Exception as e:
            logger.error(f"Error adding bullet list section: {e}")

    def _add_highlight_box_section(self, slide, content: Dict[str, Any], left, top, width, height, theme: str):
        """
        Add a highlight box section to the slide
        
        Creates a colored rectangular shape with text content to emphasize
        key points, statistics, or important information. The background color
        adapts to the selected theme while maintaining text readability.
        """
        try:
            # Create rectangular shape for the highlight box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            # Apply theme-appropriate background color for the highlight box
            fill = shape.fill
            fill.solid()
            if theme == "minimal":
                fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray background
            else:
                fill.fore_color.rgb = RGBColor(52, 152, 219)   # Blue highlight color
            
            # Add text content to the highlight box
            tf = shape.text_frame
            section_content = content.get("content", "Key Insight")
            tf.text = section_content
            
            # Style text for emphasis and readability
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # Center text in highlight box
            font = p.font
            font.size = Pt(18)  # Larger text for emphasis
            font.bold = True    # Bold weight for impact
            
            # Choose contrasting text color for readability
            if theme == "minimal":
                font.color.rgb = RGBColor(51, 51, 51)    # Dark text on light background
            else:
                font.color.rgb = RGBColor(255, 255, 255) # White text on colored background
            
            logger.info(f"Added highlight box section with content: {section_content[:100]}...")
                
        except Exception as e:
            logger.error(f"Error adding highlight box section: {e}")

    def _add_quote_section(self, slide, content: Dict[str, Any], left, top, width, height, theme: str):
        """
        Add a quote section to the slide
        
        Formats AI-generated content as an emphasized quotation with proper
        quotation marks and italic styling. Quotes are center-aligned and
        styled to stand out from regular text content.
        """
        try:
            # Create text box for the quote
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Format content with quotation marks
            section_content = content.get("content", "Quote not available")
            formatted_quote = f'"{section_content}"'  # Add quotation marks
            tf.text = formatted_quote
            
            # Style the quote with distinctive formatting
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # Center-align for emphasis
            font = p.font
            font.size = Pt(20)    # Larger than body text
            font.italic = True    # Italic styling for quotes
            
            # Apply theme-appropriate text color
            if theme == "minimal":
                font.color.rgb = RGBColor(51, 51, 51)    # Dark gray
            else:
                font.color.rgb = RGBColor(255, 255, 255) # White
            
            logger.info(f"Added quote section with content: {section_content[:100]}...")
                
        except Exception as e:
            logger.error(f"Error adding quote section: {e}")
    
    def _add_image_section(self, slide, content: Dict[str, Any], left, top, width, height, theme: str):
        """
        Add an image section to the slide
        
        Currently creates a placeholder rectangle for image content.
        This can be extended in the future to support actual image insertion
        when the frontend includes image upload functionality.
        
        Note for frontend: This section type is prepared for future image
        upload features in the builder interface.
        """
        try:
            # Create placeholder shape for image content
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            # Style the placeholder with neutral gray background
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray placeholder
            
            # Add placeholder text to indicate image location
            tf = shape.text_frame
            tf.text = "Image placeholder"
            
            # Style placeholder text
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            font = p.font
            font.size = Pt(12)  # Small text for placeholder
            font.color.rgb = RGBColor(100, 100, 100)  # Medium gray text
            
        except Exception as e:
            logger.error(f"Error adding image section: {e}")
    
    def _create_error_ppt(self, output_path: str, error_message: str) -> str:
        """
        Create a minimal PPT file with error information
        
        This fallback method ensures users always receive a downloadable file
        even when presentation generation fails. The error slide provides
        useful debugging information for both users and developers.
        
        This helps maintain a good user experience by preventing complete
        failures in the frontend download flow.
        """
        try:
            # Create minimal presentation with error information
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Position error message prominently on slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(self.slide_width - 2)  # Leave margins
            height = Inches(2)
            
            # Create text box with error details
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"Error generating slide: {error_message}"
            
            # Style error text to be clearly visible
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            font = p.font
            font.size = Pt(18)  # Large enough to read easily
            font.color.rgb = RGBColor(255, 0, 0)  # Red color for error indication
            
            # Save error presentation
            prs.save(output_path)
            logger.info(f"Error PPT created: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating error PPT: {e}")
            return output_path 